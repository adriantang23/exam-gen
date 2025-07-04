"""
PDF and PPTX Parser for Academic Documents

This module provides functionality to parse PDF and PPTX files,
extracting text content page by page or slide by slide.
Designed specifically for academic content including math symbols.
Enhanced with OCR fallback and symbol correction for better math support.
"""

import os
import sys
import re
import io
from typing import List, Union, Optional, Dict
from pathlib import Path

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("Warning: PyMuPDF not available. Install it for better PDF parsing.")

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False
    print("Warning: pdfplumber not available. Install it for fallback PDF parsing.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. Install it for PPTX parsing.")

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    print("Warning: OCR not available. Install pytesseract and Pillow for OCR fallback.")


class DocumentParser:
    """
    A parser for PDF and PPTX documents that extracts text content
    page by page or slide by slide with enhanced math symbol support.
    """
    
    def __init__(self, use_ocr_fallback: bool = True, clean_symbols: bool = True):
        """
        Initialize the parser with configuration options.
        
        Args:
            use_ocr_fallback: Whether to use OCR as fallback for problematic text
            clean_symbols: Whether to apply symbol cleaning/correction
        """
        self.use_ocr_fallback = use_ocr_fallback and OCR_AVAILABLE
        self.clean_symbols = clean_symbols
        self.supported_formats = []
        
        if PYMUPDF_AVAILABLE or PDFPLUMBER_AVAILABLE:
            self.supported_formats.append('.pdf')
        if PPTX_AVAILABLE:
            self.supported_formats.append('.pptx')
        
        # LaTeX files are always supported since they're plain text
        self.supported_formats.append('.tex')
        
        # Common symbol mappings for math/logic content
        self.symbol_mappings = {
            # Korean-looking characters that are actually math symbols
            '헍헋헎햾햫헂헍': 'true-lit',
            '햿햺헅헌햾햫헂헍': 'false-lit', 
            '헂헇헍햫헂헍': 'int-lit',
            '햿헎헇햫헆헇': 'fun-app',
            '헂헿햫헆헇': 'if-app',
            '헅헊헍햫헆헇': 'let-app',
            
            # Common mathematical symbols that get misinterpreted
            '⊢': '⊢',  # turnstile (should stay as is)
            '→': '->',  # arrow (convert to ASCII)
            '∀': 'forall',
            '∃': 'exists',
            '∧': 'and',
            '∨': 'or',
            '¬': 'not',
            '≡': 'equiv',
            '≠': '!=',
            '≤': '<=',
            '≥': '>=',
            'α': 'alpha',
            'β': 'beta',
            'γ': 'gamma',
            'δ': 'delta',
            'λ': 'lambda',
            'μ': 'mu',
            'τ': 'tau',
            'σ': 'sigma',
            'Γ': 'Gamma',
            'Δ': 'Delta',
            'Σ': 'Sigma',
            
            # Remove or replace problematic Unicode ranges
            '': '',  # Remove zero-width characters
            '\u200b': '',  # Zero-width space
            '\u200c': '',  # Zero-width non-joiner
            '\u200d': '',  # Zero-width joiner
            '\ufeff': '',  # BOM
        }
    
    def parse_document(self, file_path: Union[str, Path]) -> List[str]:
        """
        Parse a document and return text content as a list where each element
        corresponds to a page (PDF) or slide (PPTX).
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of strings, each representing the text content of a page/slide
            
        Raises:
            FileNotFoundError: If the file doesn't exist
            ValueError: If the file format is not supported
            Exception: For other parsing errors
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = file_path.suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}. "
                           f"Supported formats: {self.supported_formats}")
        
        try:
            if file_extension == '.pdf':
                return self._parse_pdf(file_path)
            elif file_extension == '.pptx':
                return self._parse_pptx(file_path)
            elif file_extension == '.tex':
                return self._parse_latex(file_path)
        except Exception as e:
            raise Exception(f"Error parsing {file_path}: {str(e)}")
    
    def _parse_pdf(self, file_path: Path) -> List[str]:
        """
        Parse a PDF file and extract text from each page with enhanced methods.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            List of strings, each representing the text content of a page
        """
        pages_text = []
        
        # Try multiple extraction methods for better results
        if PYMUPDF_AVAILABLE:
            try:
                pages_text = self._parse_pdf_pymupdf_enhanced(file_path)
                if pages_text:
                    return pages_text
            except Exception as e:
                print(f"Enhanced PyMuPDF extraction failed: {e}")
        
        # Fallback to standard PyMuPDF
        if PYMUPDF_AVAILABLE:
            try:
                pages_text = self._parse_pdf_pymupdf_standard(file_path)
                if pages_text:
                    return pages_text
            except Exception as e:
                print(f"Standard PyMuPDF failed: {e}. Trying pdfplumber...")
        
        # Fallback to pdfplumber
        if PDFPLUMBER_AVAILABLE:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        cleaned_text = self._clean_text(text if text else "")
                        pages_text.append(cleaned_text)
                return pages_text
            except Exception as e:
                print(f"pdfplumber failed: {e}")
        
        # OCR fallback if enabled
        if self.use_ocr_fallback:
            try:
                return self._parse_pdf_ocr(file_path)
            except Exception as e:
                print(f"OCR fallback failed: {e}")
        
        raise Exception("All PDF parsing methods failed")
    
    def _parse_pdf_pymupdf_enhanced(self, file_path: Path) -> List[str]:
        """Enhanced PyMuPDF extraction with multiple text extraction modes."""
        pages_text = []
        doc = fitz.open(str(file_path))
        
        try:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Try different extraction methods and pick the best one
                methods = []
                
                # Method 1: Standard text extraction
                try:
                    text1 = page.get_text()
                    methods.append(('standard', text1))
                except:
                    pass
                
                # Method 2: Extract with layout preservation
                try:
                    text2 = page.get_text("text", flags=fitz.TEXT_PRESERVE_WHITESPACE)
                    methods.append(('layout', text2))
                except:
                    pass
                
                # Method 3: Extract blocks (good for structured content)
                try:
                    blocks = page.get_text("blocks")
                    text3 = "\n".join([block[4] for block in blocks if len(block) > 4])
                    methods.append(('blocks', text3))
                except:
                    pass
                
                # Method 4: Dictionary extraction (most detailed)
                try:
                    text_dict = page.get_text("dict")
                    text4 = self._extract_from_text_dict(text_dict)
                    methods.append(('dict', text4))
                except:
                    pass
                
                # Choose the best method (most readable content)
                best_text = self._choose_best_extraction(methods)
                cleaned_text = self._clean_text(best_text)
                pages_text.append(cleaned_text)
                
        finally:
            doc.close()
        
        return pages_text
    
    def _parse_pdf_pymupdf_standard(self, file_path: Path) -> List[str]:
        """Standard PyMuPDF extraction."""
        pages_text = []
        doc = fitz.open(str(file_path))
        
        try:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                cleaned_text = self._clean_text(text)
                pages_text.append(cleaned_text)
        finally:
            doc.close()
        
        return pages_text
    
    def _parse_pdf_ocr(self, file_path: Path) -> List[str]:
        """OCR-based PDF parsing as fallback."""
        if not OCR_AVAILABLE:
            raise Exception("OCR libraries not available")
        
        pages_text = []
        doc = fitz.open(str(file_path))
        
        try:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Render page as image
                mat = fitz.Matrix(2.0, 2.0)  # High resolution
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                
                # Convert to PIL Image and run OCR
                img = Image.open(io.BytesIO(img_data))
                text = pytesseract.image_to_string(img, config='--psm 6')
                cleaned_text = self._clean_text(text)
                pages_text.append(cleaned_text)
                
        finally:
            doc.close()
        
        return pages_text
    
    def _extract_from_text_dict(self, text_dict: dict) -> str:
        """Extract text from PyMuPDF's dictionary format."""
        text_parts = []
        
        for block in text_dict.get("blocks", []):
            if "lines" in block:
                for line in block["lines"]:
                    line_text = ""
                    for span in line.get("spans", []):
                        span_text = span.get("text", "")
                        font = span.get("font", "")
                        
                        # Try to detect if this is a symbol font that needs special handling
                        if self._is_symbol_font(font):
                            span_text = self._convert_symbol_font_text(span_text, font)
                        
                        line_text += span_text
                    text_parts.append(line_text)
        
        return "\n".join(text_parts)
    
    def _is_symbol_font(self, font_name: str) -> bool:
        """Check if a font is likely to contain mathematical symbols."""
        symbol_fonts = [
            'symbol', 'math', 'times-roman', 'cambria-math', 
            'latin-modern', 'computer-modern', 'stix', 'asana'
        ]
        font_lower = font_name.lower()
        return any(sf in font_lower for sf in symbol_fonts)
    
    def _convert_symbol_font_text(self, text: str, font: str) -> str:
        """Attempt to convert symbol font text to readable form."""
        # This is a simplified approach - in practice, you'd need font-specific mappings
        # For now, we'll just apply our general symbol cleaning
        return self._clean_text(text)
    
    def _choose_best_extraction(self, methods: List[tuple]) -> str:
        """Choose the best text extraction from multiple methods."""
        if not methods:
            return ""
        
        # Score each method
        scored_methods = []
        for method_name, text in methods:
            score = self._score_text_quality(text)
            scored_methods.append((score, text, method_name))
        
        # Return the highest scoring text
        scored_methods.sort(reverse=True)
        return scored_methods[0][1]
    
    def _score_text_quality(self, text: str) -> float:
        """Score text quality based on various factors."""
        if not text or not text.strip():
            return 0.0
        
        score = 0.0
        
        # Prefer text with more readable characters
        readable_chars = sum(1 for c in text if c.isalnum() or c.isspace() or c in '.,!?-()[]{}')
        total_chars = len(text)
        if total_chars > 0:
            score += (readable_chars / total_chars) * 10
        
        # Penalize text with lots of misinterpreted symbols
        problematic_chars = sum(1 for c in text if ord(c) > 0x3000 and ord(c) < 0xD7AF)  # CJK range
        if total_chars > 0:
            score -= (problematic_chars / total_chars) * 5
        
        # Prefer text with reasonable length
        if len(text.strip()) > 10:
            score += 1
        
        # Prefer text with common English/math patterns
        if re.search(r'\b(let|if|then|else|fun|type|int|bool|string)\b', text, re.IGNORECASE):
            score += 2
        
        return score
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text."""
        if not text:
            return ""
        
        if not self.clean_symbols:
            return text.strip()
        
        # Apply symbol mappings
        for old_symbol, new_symbol in self.symbol_mappings.items():
            text = text.replace(old_symbol, new_symbol)
        
        # Remove problematic Unicode characters that often indicate misinterpretation
        # Remove characters in Korean Hangul range that are likely misinterpreted math symbols
        text = re.sub(r'[\u1100-\u11FF\u3130-\u318F\uAC00-\uD7AF]', '[SYMBOL]', text)
        
        # Clean up excessive whitespace
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)  # Max 2 consecutive newlines
        text = re.sub(r'[ \t]+', ' ', text)  # Normalize spaces
        
        # Remove zero-width characters
        text = re.sub(r'[\u200b-\u200f\ufeff]', '', text)
        
        return text.strip()
    
    def _parse_pptx(self, file_path: Path) -> List[str]:
        """
        Parse a PPTX file and extract text from each slide with comprehensive content extraction.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            List of strings, each representing the text content of a slide
        """
        if not PPTX_AVAILABLE:
            raise Exception("python-pptx not available for PPTX parsing")
        
        slides_text = []
        
        try:
            presentation = Presentation(str(file_path))
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # Extract text from all shapes in the slide (this handles tables, placeholders, and charts)
                slide_text.extend(self._extract_text_from_shapes(slide.shapes))
                
                # Extract slide notes (speaker notes) - this is separate from shapes
                notes_text = self._extract_slide_notes(slide)
                if notes_text:
                    slide_text.append(f"[NOTES: {notes_text}]")
                
                # Join all text from the slide and clean it
                combined_text = "\n".join(slide_text)
                cleaned_text = self._clean_text(combined_text)
                
                # Add slide number for reference if content exists
                if cleaned_text.strip():
                    slides_text.append(cleaned_text)
                else:
                    slides_text.append("")  # Keep empty slides for consistent indexing
            
            return slides_text
            
        except Exception as e:
            raise Exception(f"Error parsing PPTX: {e}")
    
    def _extract_text_from_shapes(self, shapes) -> List[str]:
        """Extract text from all shapes including nested shapes and groups."""
        text_content = []
        
        for shape in shapes:
            try:
                # Handle grouped shapes recursively
                if hasattr(shape, "shapes"):
                    text_content.extend(self._extract_text_from_shapes(shape.shapes))
                    continue  # Skip further processing for group shapes
                
                # Extract text content - use shape.text which is the primary text property
                # This automatically aggregates text from text frames, so we don't need both
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
                
                # Handle tables within shapes
                elif hasattr(shape, "table") and shape.table:
                    table_text = self._extract_text_from_table(shape.table)
                    if table_text:
                        text_content.extend(table_text)
                
                # Handle charts within shapes
                elif hasattr(shape, "chart") and shape.chart:
                    chart_text = self._extract_chart_text(shape.chart)
                    if chart_text:
                        text_content.extend(chart_text)
                
                # Handle complex shapes with embedded text as fallback
                elif hasattr(shape, 'element') and shape.element is not None:
                    embedded_text = self._extract_embedded_text(shape)
                    if embedded_text:
                        text_content.append(embedded_text)
                        
            except Exception as e:
                # Continue processing other shapes even if one fails
                print(f"Warning: Error extracting text from shape: {e}")
                continue
        
        return text_content
    
    def _extract_text_from_text_frame(self, text_frame) -> str:
        """Extract text from a text frame with paragraph-level detail."""
        if not text_frame or not hasattr(text_frame, 'paragraphs'):
            return ""
        
        paragraphs = []
        try:
            for paragraph in text_frame.paragraphs:
                para_text = []
                if hasattr(paragraph, 'runs'):
                    for run in paragraph.runs:
                        if hasattr(run, 'text') and run.text:
                            para_text.append(run.text)
                
                if para_text:
                    paragraphs.append(''.join(para_text))
                elif hasattr(paragraph, 'text') and paragraph.text:
                    paragraphs.append(paragraph.text)
        
        except Exception:
            # Fallback to simple text extraction
            if hasattr(text_frame, 'text'):
                return text_frame.text
        
        return '\n'.join(paragraphs)
    
    def _extract_text_from_table(self, table) -> List[str]:
        """Extract text from a table structure."""
        if not table:
            return []
        
        table_content = []
        try:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if hasattr(cell, 'text') and cell.text.strip():
                        row_text.append(cell.text.strip())
                    elif hasattr(cell, 'text_frame'):
                        cell_text = self._extract_text_from_text_frame(cell.text_frame)
                        if cell_text.strip():
                            row_text.append(cell_text.strip())
                
                if row_text:
                    table_content.append(" | ".join(row_text))
        
        except Exception as e:
            print(f"Warning: Error extracting table text: {e}")
        
        return table_content
    
    def _extract_chart_text(self, chart) -> List[str]:
        """Extract text from chart elements."""
        chart_content = []
        
        try:
            # Extract chart title
            if hasattr(chart, 'chart_title') and chart.chart_title:
                if hasattr(chart.chart_title, 'text_frame'):
                    title_text = self._extract_text_from_text_frame(chart.chart_title.text_frame)
                    if title_text:
                        chart_content.append(f"Chart Title: {title_text}")
            
            # Extract axis titles and labels
            if hasattr(chart, 'category_axis') and chart.category_axis:
                if hasattr(chart.category_axis, 'axis_title') and chart.category_axis.axis_title:
                    axis_text = self._extract_text_from_text_frame(chart.category_axis.axis_title.text_frame)
                    if axis_text:
                        chart_content.append(f"X-Axis: {axis_text}")
            
            if hasattr(chart, 'value_axis') and chart.value_axis:
                if hasattr(chart.value_axis, 'axis_title') and chart.value_axis.axis_title:
                    axis_text = self._extract_text_from_text_frame(chart.value_axis.axis_title.text_frame)
                    if axis_text:
                        chart_content.append(f"Y-Axis: {axis_text}")
            
            # Extract series names and data labels
            if hasattr(chart, 'series'):
                for i, series in enumerate(chart.series):
                    if hasattr(series, 'name') and series.name:
                        chart_content.append(f"Series {i+1}: {series.name}")
        
        except Exception:
            pass
        
        return chart_content
    
    def _extract_embedded_text(self, shape) -> str:
        """Extract any embedded text from complex shapes."""
        try:
            # This is a fallback method for complex embedded content
            if hasattr(shape, 'element') and shape.element is not None:
                # Try to get any text content from the XML element
                element = shape.element
                text_elements = element.xpath('.//a:t', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if text_elements:
                    return ' '.join([elem.text for elem in text_elements if elem.text])
        except Exception:
            pass
        
        return ""
    
    def _extract_slide_notes(self, slide) -> str:
        """Extract speaker notes from the slide."""
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_slide = slide.notes_slide
                if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                    notes_text = self._extract_text_from_text_frame(notes_slide.notes_text_frame)
                    return notes_text.strip()
        except Exception:
            pass
        
        return ""
    
    def get_document_info(self, file_path: Union[str, Path]) -> dict:
        """
        Get basic information about the document.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary with document information
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        info = {
            'file_name': file_path.name,
            'file_size': file_path.stat().st_size,
            'file_extension': file_path.suffix.lower(),
            'num_pages': 0,
            'extraction_method': 'unknown'
        }
        
        try:
            content = self.parse_document(file_path)
            info['num_pages'] = len(content)
            info['total_characters'] = sum(len(page) for page in content)
            info['non_empty_pages'] = len([page for page in content if page.strip()])
        except Exception as e:
            info['error'] = str(e)
        
        return info
    
    def _parse_latex(self, file_path: Path) -> List[str]:
        """
        Parse a LaTeX file and extract content segmented by questions/problems.
        
        Args:
            file_path: Path to the LaTeX file
            
        Returns:
            List of strings, each representing a question/problem section in raw LaTeX
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        except UnicodeDecodeError:
            # Fallback to latin-1 encoding if utf-8 fails
            with open(file_path, 'r', encoding='latin-1') as file:
                content = file.read()
        
        # Split content into logical sections
        sections = self._segment_latex_by_questions(content)
        
        # Apply minimal cleaning (just normalize whitespace, keep LaTeX intact)
        cleaned_sections = []
        for section in sections:
            # Only apply basic whitespace normalization, keep LaTeX commands
            cleaned_section = self._minimal_latex_clean(section)
            cleaned_sections.append(cleaned_section)
        
        return cleaned_sections
    
    def _segment_latex_by_questions(self, content: str) -> List[str]:
        """
        Segment LaTeX content by problems/questions using multiple detection patterns.
        Tries patterns in order of preference until one works.
        
        Args:
            content: Raw LaTeX content
            
        Returns:
            List of content sections
        """
        sections = []
        
        # Find the document content (between \begin{document} and \end{document})
        doc_start = content.find('\\begin{document}')
        doc_end = content.find('\\end{document}')
        
        if doc_start == -1:
            # No document environment found, use entire content
            document_content = content
        else:
            doc_start += len('\\begin{document}')
            if doc_end == -1:
                document_content = content[doc_start:]
            else:
                document_content = content[doc_start:doc_end]
        
        import re
        
        # Define segmentation patterns in order of preference
        segmentation_patterns = [
            # Pattern 1: \question{...} commands (highest priority)
            {
                'name': 'question_commands',
                'pattern': r'\\question\{([^}]+)\}',
                'type': 'command'
            },
            
            # Pattern 2: \begin{problem}...\end{problem} environments
            {
                'name': 'problem_environments', 
                'pattern': r'\\begin\{problem\}',
                'end_pattern': r'\\end\{problem\}',
                'type': 'environment'
            },
            
            # Pattern 3: \begin{xproblem}...\end{xproblem} environments (extra problems)
            {
                'name': 'xproblem_environments',
                'pattern': r'\\begin\{xproblem\}',
                'end_pattern': r'\\end\{xproblem\}',
                'type': 'environment'
            },
            
            # Pattern 4: Section-based divisions
            {
                'name': 'sections',
                'pattern': r'\\section\*?\{([^}]+)\}',
                'type': 'command'
            },
            
            # Pattern 5: Subsection-based divisions
            {
                'name': 'subsections',
                'pattern': r'\\subsection\*?\{([^}]+)\}',
                'type': 'command'
            },
            
            # Pattern 6: Other problem-like patterns
            {
                'name': 'paragraph_problems',
                'pattern': r'\\paragraph\{([^}]*[Pp]roblem[^}]*)\}',
                'type': 'command'
            }
        ]
        
        # Try each pattern until we find one that works
        for pattern_info in segmentation_patterns:
            sections = self._try_segmentation_pattern(document_content, pattern_info)
            if len(sections) > 1:  # Found multiple sections
                print(f"Successfully segmented using {pattern_info['name']}")
                return sections
        
        # If no patterns worked, try page breaks as delimiter
        sections = self._try_page_break_segmentation(document_content)
        if len(sections) > 1:
            print("Successfully segmented using page breaks")
            return sections
        
        # Final fallback: return entire document as one section
        print("No clear segmentation found, treating as single section")
        cleaned_content = document_content.strip()
        if cleaned_content:
            return [cleaned_content]
        else:
            return [""]
    
    def _try_segmentation_pattern(self, content: str, pattern_info: dict) -> List[str]:
        """
        Try to segment content using a specific pattern.
        
        Args:
            content: Document content
            pattern_info: Pattern information dictionary
            
        Returns:
            List of sections (empty if pattern doesn't work)
        """
        import re
        sections = []
        
        if pattern_info['type'] == 'command':
            # Handle command-based patterns like \question{...}, \section{...}
            matches = list(re.finditer(pattern_info['pattern'], content))
            
            if not matches:
                return []
            
            # Extract preamble if it contains meaningful content
            first_match_start = matches[0].start()
            preamble = content[:first_match_start].strip()
            preamble_content = self._extract_meaningful_preamble(preamble)
            if preamble_content:
                sections.append(preamble_content)
            
            # Extract each section
            for i, match in enumerate(matches):
                section_start = match.start()
                
                # Find the end of this section (start of next section or end of document)
                if i + 1 < len(matches):
                    section_end = matches[i + 1].start()
                else:
                    section_end = len(content)
                
                # Extract section content
                section_content = content[section_start:section_end].strip()
                if section_content:
                    sections.append(section_content)
        
        elif pattern_info['type'] == 'environment':
            # Handle environment-based patterns like \begin{problem}...\end{problem}
            begin_pattern = pattern_info['pattern']
            end_pattern = pattern_info['end_pattern']
            
            # Find all begin/end pairs
            begin_matches = list(re.finditer(begin_pattern, content))
            if not begin_matches:
                return []
            
            # Extract preamble
            first_begin = begin_matches[0].start()
            preamble = content[:first_begin].strip()
            preamble_content = self._extract_meaningful_preamble(preamble)
            if preamble_content:
                sections.append(preamble_content)
            
            # Extract each environment
            for begin_match in begin_matches:
                begin_pos = begin_match.start()
                
                # Find the corresponding \end{...} after this \begin{...}
                end_search_start = begin_match.end()
                end_match = re.search(end_pattern, content[end_search_start:])
                
                if end_match:
                    end_pos = end_search_start + end_match.end()
                    environment_content = content[begin_pos:end_pos].strip()
                    if environment_content:
                        sections.append(environment_content)
                else:
                    # No matching end found, take until end of document
                    environment_content = content[begin_pos:].strip()
                    if environment_content:
                        sections.append(environment_content)
        
        return sections
    
    def _try_page_break_segmentation(self, content: str) -> List[str]:
        """
        Try to segment content by page breaks.
        
        Args:
            content: Document content
            
        Returns:
            List of sections
        """
        import re
        sections = []
        
        page_break_pattern = r'\\newpage|\\clearpage'
        page_breaks = list(re.finditer(page_break_pattern, content))
        
        if not page_breaks:
            return []
        
        # Split by page breaks
        last_end = 0
        for match in page_breaks:
            section_content = content[last_end:match.start()].strip()
            if section_content:
                sections.append(section_content)
            last_end = match.end()
        
        # Add final section
        final_section = content[last_end:].strip()
        if final_section:
            sections.append(final_section)
        
        return sections
    
    def _extract_meaningful_preamble(self, preamble: str) -> str:
        """
        Extract meaningful content from document preamble, filtering out setup text.
        
        Args:
            preamble: Raw preamble content
            
        Returns:
            Meaningful preamble content or empty string
        """
        # Remove common document setup elements
        setup_patterns = [
            r'\\documentclass.*?\n',
            r'\\usepackage.*?\n',
            r'\\newcommand.*?\n',
            r'\\definecolor.*?\n',
            r'\\pagestyle.*?\n',
            r'\\fancyhf.*?\n',
            r'\\DeclareMathOperator.*?\n',
            r'\\renewcommand.*?\n',
            r'%.*?\n',  # Comments
        ]
        
        import re
        meaningful_content = preamble
        for pattern in setup_patterns:
            meaningful_content = re.sub(pattern, '', meaningful_content, flags=re.MULTILINE)
        
        # Look for actual content like titles, disclaimers, etc.
        meaningful_content = meaningful_content.strip()
        
        # If the remaining content has substantial text (not just whitespace/braces), keep it
        import re
        text_content = re.sub(r'[{}\\]', '', meaningful_content)
        text_content = re.sub(r'\s+', ' ', text_content).strip()
        
        # Keep if it has at least 20 characters of meaningful text
        if len(text_content) > 20:
            return meaningful_content
        else:
            return ""
    
    def _minimal_latex_clean(self, text: str) -> str:
        """
        Apply minimal cleaning to LaTeX text, preserving LaTeX commands and math.
        
        Args:
            text: Raw LaTeX text
            
        Returns:
            Minimally cleaned LaTeX text
        """
        if not text:
            return ""
        
        import re
        
        # Only normalize excessive whitespace, preserve LaTeX structure
        cleaned = text
        
        # Normalize excessive newlines (max 2 consecutive)
        cleaned = re.sub(r'\n\s*\n\s*\n+', '\n\n', cleaned)
        
        # Normalize spaces and tabs (but preserve intentional spacing)
        cleaned = re.sub(r'[ \t]+', ' ', cleaned)
        
        # Trim leading/trailing whitespace from each line while preserving indentation structure
        lines = cleaned.split('\n')
        normalized_lines = []
        for line in lines:
            # Only strip trailing whitespace, preserve leading whitespace for LaTeX indentation
            normalized_lines.append(line.rstrip())
        
        cleaned = '\n'.join(normalized_lines)
        
        return cleaned.strip()


def parse_document(file_path: Union[str, Path], use_ocr_fallback: bool = True, 
                  clean_symbols: bool = True) -> List[str]:
    """
    Convenience function to parse a document with enhanced options.
    
    Args:
        file_path: Path to the document file
        use_ocr_fallback: Whether to use OCR as fallback for problematic text
        clean_symbols: Whether to apply symbol cleaning/correction
        
    Returns:
        List of strings, each representing the text content of a page/slide
    """
    parser = DocumentParser(use_ocr_fallback=use_ocr_fallback, clean_symbols=clean_symbols)
    return parser.parse_document(file_path)


def main():
    """
    Command line interface for testing the parser.
    """
    if len(sys.argv) < 2:
        print("Usage: python parser.py <file_path> [--no-ocr] [--no-clean]")
        print("Examples:")
        print("  python parser.py scanable_pdf_test_documents/02-the-basics-A1.pdf")
        print("  python parser.py scanable_pdf_test_documents/Reinforcement_Learning.pptx")
        print("  python parser.py scanable_pdf_test_documents/wa3.tex")
        print("  python parser.py scanable_pdf_test_documents/CS237hw5.tex")
        print()
        print("Supported formats:")
        print("  PDF  - Multiple extraction methods with OCR fallback")
        print("  PPTX - Comprehensive content extraction (slides, notes, charts)")
        print("  LaTeX - Auto-detects structure: \\question{}, \\begin{problem}, \\section{}, etc.")
        print()
        print("Options:")
        print("  --no-ocr    Disable OCR fallback (PDF only)")
        print("  --no-clean  Disable symbol cleaning")
        return
    
    file_path = sys.argv[1]
    use_ocr = '--no-ocr' not in sys.argv
    clean_symbols = '--no-clean' not in sys.argv
    
    try:
        parser = DocumentParser(use_ocr_fallback=use_ocr, clean_symbols=clean_symbols)
        
        # Get document info
        info = parser.get_document_info(file_path)
        print(f"Document Info:")
        print(f"  Name: {info['file_name']}")
        print(f"  Size: {info['file_size']} bytes")
        print(f"  Type: {info['file_extension']}")
        print(f"  Pages/Slides: {info.get('num_pages', 'Unknown')}")
        if 'total_characters' in info:
            print(f"  Total Characters: {info['total_characters']}")
        if 'non_empty_pages' in info:
            print(f"  Non-empty: {info['non_empty_pages']}")
        print(f"  OCR Fallback: {'Enabled' if use_ocr else 'Disabled'}")
        print(f"  Symbol Cleaning: {'Enabled' if clean_symbols else 'Disabled'}")
        print()
        
        # Parse the document
        pages = parser.parse_document(file_path)
        
        # Determine appropriate terminology based on file type
        if info['file_extension'] == '.pdf':
            unit_name = 'pages'
        elif info['file_extension'] == '.pptx':
            unit_name = 'slides'
        elif info['file_extension'] == '.tex':
            unit_name = 'sections'
        else:
            unit_name = 'sections'
        
        print(f"Successfully parsed {len(pages)} {unit_name}:")
        print("=" * 50)
        
        for i, page_content in enumerate(pages, 1):
            unit_label = unit_name[:-1].title()  # Convert 'pages' to 'Page', etc.
            print(f"\n--- {unit_label} {i} ---")
            if page_content:
                # Show first 200 characters of each page
                preview = page_content[:200]
                if len(page_content) > 200:
                    preview += "..."
                print(preview)
            else:
                print("(No text content)")
            print(f"[Length: {len(page_content)} characters]")
        
    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    main()
